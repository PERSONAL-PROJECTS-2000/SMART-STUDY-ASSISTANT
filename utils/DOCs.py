import io
import os
from PIL import Image
from PyPDF2 import PdfReader
from pptx import Presentation
import docx
import openpyxl
import pytesseract
from pdf2image import convert_from_bytes

def extract_text_from_file(file):
    if file is None:
        return ""
    file_ext = os.path.splitext(file.name)[1].lower()
    try:
        if file_ext == '.txt':
            return file.read().decode('utf-8')
        elif file_ext == '.pdf':
            text = ""
            try:
                reader = PdfReader(io.BytesIO(file.read()))
                for page in reader.pages:
                    text += page.extract_text() + "\\n"
            except:
                pass
            if len(text.strip()) < 50:
                file.seek(0)
                images = convert_from_bytes(file.read())
                for img in images:
                    text += pytesseract.image_to_string(img) + "\\n"
            return text
        elif file_ext == '.docx' or file_ext == '.doc':
            doc = docx.Document(io.BytesIO(file.read()))
            return "\\n".join([para.text for para in doc.paragraphs])
        elif file_ext == '.pptx':
            prs = Presentation(io.BytesIO(file.read()))
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\\n"
            return text
        elif file_ext == '.xlsx':
            wb = openpyxl.load_workbook(io.BytesIO(file.read()))
            text = ""
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    text += " ".join([str(cell) for cell in row if cell]) + "\\n"
            return text
        elif file_ext in ['.jpg', '.jpeg', '.png']:
            img = Image.open(io.BytesIO(file.read()))
            return pytesseract.image_to_string(img)
    except Exception as e:
        return f"Error extracting text: {str(e)}"
    return ""